from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw
import os

# Constants
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR = os.path.join(SCRIPT_DIR, '../assets')
OUTPUT_FILE = os.path.join(SCRIPT_DIR, 'pyconsg2026_template.pptx')

# Colors
DARK_BG_RGB = (10, 10, 20)
PY_YELLOW_RGB = (255, 212, 59)
SG_RED_RGB = (239, 51, 64)
SG_ORCHID_RGB = (153, 50, 204)
WHITE_RGB = (255, 255, 255)

def create_gradient_overlay(width, height, color_rgb, start_alpha, end_alpha):
    """Creates a vertical gradient image from start_alpha to end_alpha of a given color."""
    base = Image.new('RGBA', (width, height), (0, 0, 0, 0))
    top = Image.new('RGBA', (width, height), color_rgb + (0,))
    
    # Create alpha mask
    mask = Image.new('L', (width, height))
    draw = ImageDraw.Draw(mask)
    
    for y in range(height):
        alpha = int(start_alpha + (end_alpha - start_alpha) * (y / height))
        draw.line((0, y, width, y), fill=alpha)
    
    top.putalpha(mask)
    return top

def add_title_slide(prs):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg_path = os.path.join(ASSETS_DIR, 'hero-bg.png')
    if os.path.exists(bg_path):
        left = top = 0
        slide.shapes.add_picture(bg_path, left, top, width=prs.slide_width, height=prs.slide_height)
    
    # Gradient Overlay
    # Generate overlay image
    overlay_path = os.path.join(SCRIPT_DIR, 'overlay.png')
    # 70% opacity = 178, 100% opacity = 255
    overlay_img = create_gradient_overlay(int(prs.slide_width.pt), int(prs.slide_height.pt), DARK_BG_RGB, 178, 255)
    overlay_img.save(overlay_path)
    slide.shapes.add_picture(overlay_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    os.remove(overlay_path)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Python in the Lion City"
    p.font.name = 'Arial' # Fallback for Outfit
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE_RGB)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = "PyCon SG 2026"
    p.font.name = 'Courier New' # Fallback for Space Grotesk
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(*PY_YELLOW_RGB)
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(1))
    tf = date_box.text_frame
    p = tf.add_paragraph()
    p.text = "19th - 21st June 2026"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*SG_RED_RGB)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Solid Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*DARK_BG_RGB)
    
    # Header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = "Content Slide Title"
    p.font.name = 'Arial'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*PY_YELLOW_RGB)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "• First point goes here"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*WHITE_RGB)
    p.space_after = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Second point with details"
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*WHITE_RGB)
    p.space_after = Pt(14)

def add_photo_slide(prs, title, image_filename, is_dark=True):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background Image
    bg_path = os.path.join(ASSETS_DIR, image_filename)
    if os.path.exists(bg_path):
        slide.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # Overlay for text readability
    # Dark overlay for both modes to ensure white/yellow text pops, 
    # or light overlay for light mode? 
    # Let's stick to the dark theme aesthetic of the conference for now, 
    # but maybe lighter for 'Day' shots if we want to keep them bright.
    # However, user asked for 'Light Mode' slides.
    
    if is_dark:
        # Dark mode: Dark overlay (similar to title slide)
        overlay_color = DARK_BG_RGB
        text_color = WHITE_RGB
        accent_color = PY_YELLOW_RGB
        alpha_start, alpha_end = 100, 200 # Lighter overlay than title
    else:
        # Light mode: Maybe a white fade?
        overlay_color = WHITE_RGB
        text_color = DARK_BG_RGB
        accent_color = SG_RED_RGB
        alpha_start, alpha_end = 150, 230
        
    overlay_path = os.path.join(SCRIPT_DIR, f'overlay_{image_filename}.png')
    overlay_img = create_gradient_overlay(int(prs.slide_width.pt), int(prs.slide_height.pt), overlay_color, alpha_start, alpha_end)
    overlay_img.save(overlay_path)
    slide.shapes.add_picture(overlay_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
    os.remove(overlay_path)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*text_color)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle/Accent
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.add_paragraph()
    p.text = "PyCon SG 2026"
    p.font.name = 'Courier New'
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(*accent_color)
    p.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs)
    add_content_slide(prs)
    
    # Additional Background Slides
    locations = [
        ("Changi Jewel (Day)", "changijewelday.jpeg", False),
        ("Changi Jewel (Night)", "changijewelnight.jpeg", True),
        ("Gardens by the Bay (Day)", "gardensbythebayday.jpeg", False),
        ("Gardens by the Bay (Night)", "gardensbythebaynight.jpeg", True),
        ("Marina Bay Sands (Day)", "marinabaysandsday.jpeg", False),
        ("Marina Bay Sands (Night)", "marinabaysandsnight.jpeg", True),
    ]
    
    for title, filename, is_dark in locations:
        add_photo_slide(prs, title, filename, is_dark)
    
    # Save
    prs.save(OUTPUT_FILE)
    print(f"Created {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
